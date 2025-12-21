"""
Script pour générer des fichiers de test (PDF, DOCX, XLSX, images)
Utilisé par les tests pour créer des fixtures dynamiques
"""
import os
from pathlib import Path
from typing import Optional

# Répertoires
FIXTURES_DIR = Path(__file__).parent
DOCS_DIR = FIXTURES_DIR / "documents"
IMAGES_DIR = FIXTURES_DIR / "images"


def generate_pdf(output_path: Optional[Path] = None) -> Path:
    """Génère un fichier PDF de test avec reportlab"""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
    except ImportError:
        print("WARNING: reportlab not installed. Skipping PDF generation.")
        print("Install with: pip install reportlab")
        return None

    if output_path is None:
        output_path = DOCS_DIR / "sample.pdf"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Créer le document
    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Titre
    title = Paragraph("<b>Document de Test PDF</b>", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 0.2 * inch))

    # Contenu page 1
    content1 = Paragraph("""
    Ce document PDF est généré automatiquement pour tester le système d'ingestion v2.0.
    Il contient plusieurs pages avec du texte structuré.
    """, styles['Normal'])
    story.append(content1)
    story.append(Spacer(1, 0.2 * inch))

    # Liste
    for i in range(1, 6):
        item = Paragraph(f"• Point {i}: Information de test numéro {i}", styles['Normal'])
        story.append(item)

    story.append(PageBreak())

    # Page 2
    title2 = Paragraph("<b>Page 2: Chunking Sémantique</b>", styles['Heading1'])
    story.append(title2)
    story.append(Spacer(1, 0.2 * inch))

    content2 = Paragraph("""
    Le chunking sémantique respecte la structure du document.
    Il découpe le texte en segments cohérents plutôt qu'en blocs arbitraires de caractères.
    Cela améliore la qualité de la recherche et la pertinence des résultats RAG.
    """, styles['Normal'])
    story.append(content2)

    # Construire le PDF
    doc.build(story)
    print(f"✅ PDF créé: {output_path}")
    return output_path


def generate_docx(output_path: Optional[Path] = None) -> Path:
    """Génère un fichier DOCX de test avec python-docx"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ImportError:
        print("WARNING: python-docx not installed. Skipping DOCX generation.")
        print("Install with: pip install python-docx")
        return None

    if output_path is None:
        output_path = DOCS_DIR / "sample.docx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Créer le document
    doc = Document()

    # Titre
    doc.add_heading('Document DOCX de Test', 0)

    # Paragraphe
    doc.add_paragraph('Ce document Word teste l\'extraction avec Unstructured.io.')

    # Liste
    doc.add_heading('Liste de fonctionnalités v2.0:', level=2)
    features = [
        '13 formats de fichiers supportés',
        'Chunking sémantique avec LangChain',
        'Déduplication automatique par hash SHA256',
        'Métadonnées enrichies (11 champs)',
        'OCR intégré avec Tesseract'
    ]
    for feature in features:
        doc.add_paragraph(feature, style='List Bullet')

    # Tableau
    doc.add_heading('Tableau de test:', level=2)
    table = doc.add_table(rows=4, cols=3)
    table.style = 'Light Grid Accent 1'

    # En-têtes
    headers = ['Format', 'Outil', 'Support OCR']
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header

    # Données
    data = [
        ['PDF', 'Unstructured', 'Oui'],
        ['DOCX', 'Unstructured', 'Non'],
        ['Images', 'Tesseract', 'Oui']
    ]
    for i, row in enumerate(data, start=1):
        for j, value in enumerate(row):
            table.rows[i].cells[j].text = value

    # Sauvegarder
    doc.save(str(output_path))
    print(f"✅ DOCX créé: {output_path}")
    return output_path


def generate_xlsx(output_path: Optional[Path] = None) -> Path:
    """Génère un fichier XLSX de test avec openpyxl"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
    except ImportError:
        print("WARNING: openpyxl not installed. Skipping XLSX generation.")
        print("Install with: pip install openpyxl")
        return None

    if output_path is None:
        output_path = DOCS_DIR / "sample.xlsx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Créer le workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Données Test"

    # En-têtes avec style
    headers = ['Technologie', 'Catégorie', 'Année', 'Description']
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")

    # Données
    data = [
        ['Python', 'Langage', 1991, 'Langage polyvalent et lisible'],
        ['FastAPI', 'Framework Web', 2018, 'Framework moderne et performant'],
        ['ChromaDB', 'Base vectorielle', 2022, 'Stockage d\'embeddings'],
        ['Ollama', 'Serveur LLM', 2023, 'Exécution locale de modèles'],
        ['LangChain', 'Framework', 2022, 'Développement d\'applications LLM']
    ]

    for row_idx, row_data in enumerate(data, start=2):
        for col_idx, value in enumerate(row_data, start=1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    # Ajuster les largeurs de colonnes
    ws.column_dimensions['A'].width = 15
    ws.column_dimensions['B'].width = 20
    ws.column_dimensions['C'].width = 10
    ws.column_dimensions['D'].width = 40

    # Sauvegarder
    wb.save(str(output_path))
    print(f"✅ XLSX créé: {output_path}")
    return output_path


def generate_pptx(output_path: Optional[Path] = None) -> Path:
    """Génère un fichier PPTX de test avec python-pptx"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("WARNING: python-pptx not installed. Skipping PPTX generation.")
        print("Install with: pip install python-pptx")
        return None

    if output_path is None:
        output_path = DOCS_DIR / "sample.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Titre
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Système d'Ingestion v2.0"
    subtitle.text = "Test de l'extraction PowerPoint"

    # Slide 2: Contenu avec liste
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    body2 = slide2.placeholders[1]
    title2.text = "Nouvelles Fonctionnalités"

    tf = body2.text_frame
    tf.text = "Support multi-formats"

    bullet_points = [
        "13 formats de fichiers",
        "Chunking sémantique intelligent",
        "Déduplication SHA256",
        "Métadonnées enrichies",
        "OCR avec Tesseract"
    ]

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1

    # Sauvegarder
    prs.save(str(output_path))
    print(f"✅ PPTX créé: {output_path}")
    return output_path


def generate_image_with_text(output_path: Optional[Path] = None) -> Path:
    """Génère une image PNG avec du texte (pour tester OCR)"""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("WARNING: Pillow not installed. Skipping image generation.")
        print("Install with: pip install Pillow")
        return None

    if output_path is None:
        output_path = IMAGES_DIR / "sample_ocr.png"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Créer une image blanche
    img = Image.new('RGB', (800, 400), color='white')
    draw = ImageDraw.Draw(img)

    # Essayer d'utiliser une police système, sinon utiliser la police par défaut
    try:
        font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
        font_body = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
    except:
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()

    # Texte
    draw.text((50, 50), "Test OCR - Ingestion v2.0", fill='black', font=font_title)
    draw.text((50, 120), "Cette image contient du texte", fill='black', font=font_body)
    draw.text((50, 160), "qui doit être extrait par Tesseract.", fill='black', font=font_body)
    draw.text((50, 200), "Le système v2 supporte l'OCR", fill='black', font=font_body)
    draw.text((50, 240), "pour PDF scannés et images.", fill='black', font=font_body)

    # Bordure
    draw.rectangle([(10, 10), (790, 390)], outline='blue', width=3)

    # Sauvegarder
    img.save(str(output_path))
    print(f"✅ Image PNG créée: {output_path}")
    return output_path


def generate_all():
    """Génère tous les fichiers de test"""
    print("\n🔨 Génération des fichiers de test...")
    print("=" * 60)

    # Créer les répertoires
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)

    # Générer les fichiers
    generate_pdf()
    generate_docx()
    generate_xlsx()
    generate_pptx()
    generate_image_with_text()

    print("=" * 60)
    print("✅ Génération terminée !\n")


if __name__ == "__main__":
    generate_all()
